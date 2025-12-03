from jinja2 import Environment, FileSystemLoader
import pandas as pd
import os
import re
import markdown
from bs4 import BeautifulSoup
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, ListFlowable, ListItem
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from io import BytesIO

def convert_markdown_to_reportlab(md_text: str, styles):
    """
    Convert markdown text to ReportLab flowables using markdown library.
    Handles headings, bold, italic, lists, tables, and paragraphs.
    """
    # Convert markdown to HTML
    html = markdown.markdown(md_text, extensions=['tables', 'nl2br'])
    
    # Parse HTML with BeautifulSoup
    soup = BeautifulSoup(html, 'html.parser')
    
    elements = []
    
    # Define custom styles
    h1_style = ParagraphStyle(
        'CustomH1',
        parent=styles['Heading1'],
        fontSize=18,
        textColor=colors.HexColor('#2C3E50'),
        spaceAfter=12,
        spaceBefore=16,
        fontName='Helvetica-Bold'
    )
    
    h2_style = ParagraphStyle(
        'CustomH2',
        parent=styles['Heading2'],
        fontSize=15,
        textColor=colors.HexColor('#34495E'),
        spaceAfter=10,
        spaceBefore=14,
        fontName='Helvetica-Bold'
    )
    
    h3_style = ParagraphStyle(
        'CustomH3',
        parent=styles['Heading3'],
        fontSize=13,
        textColor=colors.HexColor('#34495E'),
        spaceAfter=8,
        spaceBefore=12,
        fontName='Helvetica-Bold'
    )
    
    body_style = ParagraphStyle(
        'BodyText',
        parent=styles['Normal'],
        fontSize=11,
        leading=16,
        alignment=TA_JUSTIFY,
        spaceAfter=10
    )
    
    list_style = ParagraphStyle(
        'ListItem',
        parent=styles['Normal'],
        fontSize=11,
        leading=16,
        leftIndent=20,
        bulletIndent=10,
        spaceAfter=6
    )
    
    # Process each HTML element
    for element in soup.find_all(recursive=False):
        if element.name == 'h1':
            elements.append(Paragraph(clean_html_tags(str(element)), h1_style))
        elif element.name == 'h2':
            elements.append(Paragraph(clean_html_tags(str(element)), h2_style))
        elif element.name == 'h3':
            elements.append(Paragraph(clean_html_tags(str(element)), h3_style))
        elif element.name == 'p':
            text = clean_html_tags(str(element))
            if text.strip():
                elements.append(Paragraph(text, body_style))
        elif element.name in ['ul', 'ol']:
            # Handle lists
            for li in element.find_all('li', recursive=False):
                text = clean_html_tags(str(li))
                elements.append(Paragraph(f"• {text}", list_style))
        elif element.name == 'table':
            # Handle tables
            table_data = []
            for row in element.find_all('tr'):
                row_data = []
                for cell in row.find_all(['th', 'td']):
                    row_data.append(clean_html_tags(str(cell)))
                if row_data:
                    table_data.append(row_data)
            
            if table_data:
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3498DB')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ]))
                elements.append(t)
                elements.append(Spacer(1, 0.15*inch))
        elif element.name == 'blockquote':
            quote_style = ParagraphStyle(
                'Quote',
                parent=styles['Normal'],
                fontSize=11,
                leftIndent=30,
                rightIndent=30,
                textColor=colors.HexColor('#7F8C8D'),
                spaceAfter=10
            )
            text = clean_html_tags(str(element))
            elements.append(Paragraph(text, quote_style))
    
    return elements

def clean_html_tags(html_text: str) -> str:
    """
    Convert HTML tags to ReportLab XML tags.
    """
    # Remove the outer tag
    html_text = re.sub(r'^<[^>]+>', '', html_text)
    html_text = re.sub(r'</[^>]+>$', '', html_text)
    
    # Convert <strong> and <b> to <b>
    html_text = re.sub(r'<strong>', '<b>', html_text)
    html_text = re.sub(r'</strong>', '</b>', html_text)
    
    # Convert <em> and <i> to <i>
    html_text = re.sub(r'<em>', '<i>', html_text)
    html_text = re.sub(r'</em>', '</i>', html_text)
    
    # Convert <code> to monospace
    html_text = re.sub(r'<code>', '<font name="Courier" size="10">', html_text)
    html_text = re.sub(r'</code>', '</font>', html_text)
    
    # Remove other HTML tags
    html_text = re.sub(r'<br\s*/?>', '<br/>', html_text)
    
    # Remove any remaining unsupported tags
    html_text = re.sub(r'<(?!/?[biu]|/?font|br/?)([^>]+)>', '', html_text)
    
    return html_text.strip()

def parse_markdown_to_paragraphs(text: str, styles):
    """
    DEPRECATED: Use convert_markdown_to_reportlab instead.
    Kept for backward compatibility.
    """
    return convert_markdown_to_reportlab(text, styles)

def render_html_report(context: dict, template_dir: str, template_file: str) -> str:
    env = Environment(loader=FileSystemLoader(template_dir))
    template = env.get_template(template_file)
    html = template.render(**context)
    return html

def save_pdf_from_html(html_str: str, output_path: str):
    """
    This function is kept for compatibility but now uses context-based PDF generation.
    For better results, use save_pdf_from_context() directly.
    """
    # Simple fallback - just save as text
    from reportlab.platypus import SimpleDocTemplate, Paragraph
    from reportlab.lib.styles import getSampleStyleSheet
    
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Strip HTML tags for simple rendering
    import re
    text = re.sub('<[^<]+?>', '', html_str)
    story.append(Paragraph(text[:1000], styles['Normal']))  # Limit to first 1000 chars
    
    doc.build(story)
    print(f"Saved PDF report to {output_path}")

def save_pdf_from_context(context: dict, output_path: str):
    """
    Generate a professional PDF report from context using ReportLab.
    This is the recommended method for PDF generation.
    """
    doc = SimpleDocTemplate(output_path, pagesize=letter,
                           topMargin=0.5*inch, bottomMargin=0.5*inch)
    
    # Define styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#2C3E50'),
        spaceAfter=30,
        alignment=TA_CENTER
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#34495E'),
        spaceAfter=12,
        spaceBefore=12
    )
    
    body_style = ParagraphStyle(
        'BodyText',
        parent=styles['Normal'],
        fontSize=11,
        leading=16,
        alignment=TA_JUSTIFY,
        spaceAfter=8
    )
    
    # Build story
    story = []
    
    # Title
    story.append(Paragraph(context.get('title', 'TrendSpotter Report'), title_style))
    story.append(Paragraph(f"Generated: {context.get('generated_on', 'N/A')}", styles['Normal']))
    story.append(Spacer(1, 0.3*inch))
    
    # File Info
    if 'file_info' in context:
        story.append(Paragraph("File Information", heading_style))
        file_data = [
            ['Property', 'Value'],
            ['Filename', context['file_info'].get('filename', 'N/A')],
            ['Total Rows', str(context['file_info'].get('rows', 0))],
            ['Total Columns', str(context['file_info'].get('columns', 0))]
        ]
        t = Table(file_data, colWidths=[2*inch, 4*inch])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3498DB')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(t)
        story.append(Spacer(1, 0.2*inch))
    
    # Anomaly Summary
    if 'anomalies_summary' in context:
        story.append(Paragraph("Anomaly Detection Summary", heading_style))
        anomaly_data = [
            ['Metric', 'Value'],
            ['Total Anomalies', str(context['anomalies_summary'].get('total', 0))],
            ['Percentage', context['anomalies_summary'].get('percentage', '0%')]
        ]
        t = Table(anomaly_data, colWidths=[2*inch, 4*inch])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#E74C3C')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(t)
        story.append(Spacer(1, 0.2*inch))
    
    # AI Analysis
    if 'ai_analysis' in context and context['ai_analysis']:
        story.append(Paragraph("AI-Generated Business Analysis", heading_style))
        story.append(Spacer(1, 0.15*inch))
        
        # Convert markdown to formatted PDF elements
        analysis_text = context['ai_analysis']
        formatted_elements = convert_markdown_to_reportlab(analysis_text, styles)
        story.extend(formatted_elements)
    
    # Build PDF
    doc.build(story)
    print(f"Saved PDF report to {output_path}")

def save_pptx_from_context(output_pptx: str, context: dict):
    """
    Generate a professional PowerPoint report from context using python-pptx.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    # Create a presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = context.get('title', 'TrendSpotter Report')
    subtitle.text = f"Generated: {context.get('generated_on', '')}"
    
    # Slide 2: File Information
    if 'file_info' in context:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'File Information'
        
        tf = body_shape.text_frame
        tf.text = f"Filename: {context['file_info'].get('filename', 'N/A')}"
        
        p = tf.add_paragraph()
        p.text = f"Total Rows: {context['file_info'].get('rows', 0)}"
        
        p = tf.add_paragraph()
        p.text = f"Total Columns: {context['file_info'].get('columns', 0)}"
        
        if 'column_names' in context['file_info']:
            p = tf.add_paragraph()
            p.text = f"Columns: {', '.join(context['file_info']['column_names'][:5])}"
    
    # Slide 3: Anomaly Summary
    if 'anomalies_summary' in context:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'Anomaly Detection Summary'
        
        tf = body_shape.text_frame
        tf.text = f"Total Anomalies: {context['anomalies_summary'].get('total', 0)}"
        
        p = tf.add_paragraph()
        p.text = f"Percentage: {context['anomalies_summary'].get('percentage', '0%')}"
    
    # Slide 4: AI Analysis
    if 'ai_analysis' in context and context['ai_analysis']:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'AI-Generated Analysis'
        
        tf = body_shape.text_frame
        # Split analysis into chunks if too long
        analysis_text = context['ai_analysis']
        if len(analysis_text) > 500:
            analysis_text = analysis_text[:500] + "..."
        tf.text = analysis_text
    
    # Save
    prs.save(output_pptx)
    print(f"Saved PPTX report to {output_pptx}")

if __name__ == "__main__":
    context = {
      "title": "Campaign Report",
      "generated_on": pd.Timestamp.now().strftime("%Y-%m-%d %H:%M:%S"),
      "summary_table_html": pd.DataFrame([[1,2,3],[4,5,6]], columns=["A","B","C"]).to_html(index=False),
      "anomalies_table_html": pd.DataFrame([[10, 'row1'], [20, 'row2']], columns=["score","row_id"]).to_html(index=False),
      "ai_analysis": "<h2>AI Summary</h2><p>Some analysis text...</p>",
      "sample_table": [["Metric","Value"], ["Spend","$1000"], ["Clicks","300"]],
    }
    html = render_html_report(context, template_dir="templates", template_file="report_template.html")
    save_pdf_from_html(html, "report.pdf")
    save_pptx_from_context("template.pptx", "report.pptx", context)